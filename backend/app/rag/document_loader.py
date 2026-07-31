from __future__ import annotations

from pathlib import Path
from typing import List, Dict, Union
import csv
import re

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

from app.core.logging import get_logger

logger = get_logger(__name__)

SUPPORTED_EXTENSIONS = {'.pdf': 'pdf', '.docx': 'docx', '.txt': 'txt', '.csv': 'csv', '.pptx': 'pptx'}


class DocumentLoader:
    """Load and normalize supported college documents (PDF, DOCX, TXT) into clean text."""

    def __init__(self, documents_dir: Path | str | None = None) -> None:
        if documents_dir:
            self.documents_dir = Path(documents_dir)
        else:
            self.documents_dir = Path(__file__).resolve().parents[2] / 'documents'

    def clean_text(self, text: str) -> str:
        """Clean and sanitize raw extracted text for optimal chunking and embedding."""
        if not text:
            return ""
        # Replace non-breaking spaces & control characters
        text = text.replace('\xa0', ' ').replace('\r\n', '\n').replace('\r', '\n')
        # Remove null bytes or non-printable ASCII/Unicode control characters except newlines/tabs
        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
        # Collapse multi-space lines and limit duplicate blank lines
        lines = [re.sub(r'[ \t]+', ' ', line).strip() for line in text.split('\n')]
        cleaned = '\n'.join(lines)
        cleaned = re.sub(r'\n{3,}', '\n\n', cleaned)
        return cleaned.strip()

    def load_single_document(self, file_path: Union[str, Path]) -> Dict[str, str]:
        """Extract and clean text from a single file path."""
        path = Path(file_path)
        if not path.exists() or not path.is_file():
            raise ValueError(f"Document file does not exist: {file_path}")

        suffix = path.suffix.lower()
        if suffix not in SUPPORTED_EXTENSIONS:
            raise ValueError(f"Unsupported file type '{suffix}'. Allowed: PDF, DOCX, TXT")

        raw_text = self._extract_text(path)
        cleaned = self.clean_text(raw_text)

        if not cleaned:
            raise ValueError(f"Could not extract readable text from '{path.name}' (file may be empty, image-only, or corrupted).")

        return {
            'source': str(path.resolve()),
            'filename': path.name,
            'file_type': suffix.lstrip('.'),
            'content': cleaned
        }

    def load_documents(self) -> List[Dict[str, str]]:
        """Read all supported files from the target directory."""
        if not self.documents_dir.exists():
            logger.warning('Documents directory does not exist: %s', self.documents_dir)
            return []

        documents: List[Dict[str, str]] = []
        for file_path in sorted(self.documents_dir.rglob('*')):
            if not file_path.is_file():
                continue

            suffix = file_path.suffix.lower()
            if suffix not in SUPPORTED_EXTENSIONS:
                continue

            try:
                doc = self.load_single_document(file_path)
                documents.append(doc)
                logger.info('Loaded document: %s', file_path)
            except Exception as exc:
                logger.warning('Failed to load document %s: %s', file_path, exc)

        logger.info('Loaded %s document(s) from %s', len(documents), self.documents_dir)
        return documents

    def _extract_text(self, file_path: Path) -> str:
        """Extract text from PDF, DOCX, TXT, PPTX or CSV files."""
        suffix = file_path.suffix.lower()
        if suffix == '.pdf':
            return self._read_pdf(file_path)
        if suffix == '.docx':
            return self._read_docx(file_path)
        if suffix == '.txt':
            return self._read_text(file_path)
        if suffix == '.csv':
            return self._read_csv(file_path)
        if suffix == '.pptx':
            return self._read_pptx(file_path)
        return ''

    def _read_pdf(self, file_path: Path) -> str:
        """Extract text from a PDF file using PyMuPDF (fitz) with pypdf fallback."""
        # 1. Try PyMuPDF (fitz)
        if fitz is not None:
            try:
                doc = fitz.open(str(file_path))
                text_pages = []
                for page in doc:
                    text_pages.append(page.get_text("text") or "")
                doc.close()
                extracted = "\n".join(text_pages).strip()
                if extracted:
                    logger.info("Successfully extracted text from %s via PyMuPDF", file_path.name)
                    return extracted
            except Exception as exc:
                logger.warning("PyMuPDF failed for %s: %s. Falling back to pypdf.", file_path, exc)

        # 2. Fallback to pypdf
        if PdfReader is not None:
            try:
                reader = PdfReader(str(file_path))
                pages = [page.extract_text() or '' for page in reader.pages]
                return '\n'.join(page for page in pages if page).strip()
            except Exception as exc:
                logger.exception('pypdf failed to read PDF: %s', file_path)
                return ''

        logger.warning('Neither PyMuPDF nor pypdf available for %s', file_path)
        return ''

    def _read_docx(self, file_path: Path) -> str:
        """Extract text from a DOCX file."""
        if DocxDocument is None:
            logger.warning('python-docx is not installed; skipping DOCX extraction for %s', file_path)
            return ''

        try:
            document = DocxDocument(str(file_path))
            paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
            return '\n'.join(paragraphs).strip()
        except Exception as exc:
            logger.exception('Failed to read DOCX: %s', file_path)
            return ''

    def _read_text(self, file_path: Path) -> str:
        """Read a plain text file using UTF-8 with fallback."""
        try:
            return file_path.read_text(encoding='utf-8').strip()
        except Exception:
            try:
                return file_path.read_text(encoding='latin-1').strip()
            except Exception as exc:
                logger.exception('Failed to read text file: %s', file_path)
                return ''

    def _read_csv(self, file_path: Path) -> str:
        """Extract text from a CSV file."""
        try:
            with open(file_path, newline='', encoding='utf-8') as f:
                reader = csv.reader(f)
                rows = [' '.join(row) for row in reader if row]
                return '\n'.join(rows).strip()
        except Exception as exc:
            logger.exception('Failed to read CSV: %s', file_path)
            return ''

    def _read_pptx(self, file_path: Path) -> str:
        """Extract text from a PPTX file."""
        if Presentation is None:
            logger.warning('python-pptx is not installed; skipping PPTX extraction for %s', file_path)
            return ''
            
        try:
            prs = Presentation(str(file_path))
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return '\n'.join(text_runs).strip()
        except Exception as exc:
            logger.exception('Failed to read PPTX: %s', file_path)
            return ''

